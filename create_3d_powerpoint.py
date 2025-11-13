#!/usr/bin/env python3
"""
Script to create a PowerPoint presentation with 3D models
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

def create_3d_presentation():
    """Create a PowerPoint presentation with 3D model placeholders and information"""
    
    # Create presentation object
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title Slide
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    
    # Add gradient background
    background = slide1.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(41, 128, 185)
    
    # Title
    title_box = slide1.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame = title_box.text_frame
    title_frame.text = "Presentación con Modelos 3D"
    title_para = title_frame.paragraphs[0]
    title_para.font.size = Pt(54)
    title_para.font.bold = True
    title_para.font.color.rgb = RGBColor(255, 255, 255)
    title_para.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box = slide1.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(0.8))
    subtitle_frame = subtitle_box.text_frame
    subtitle_frame.text = "Explorando el mundo tridimensional"
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.font.size = Pt(28)
    subtitle_para.font.color.rgb = RGBColor(236, 240, 241)
    subtitle_para.alignment = PP_ALIGN.CENTER
    
    # Slide 2: Introduction to 3D Models
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    background2 = slide2.background
    fill2 = background2.fill
    fill2.solid()
    fill2.fore_color.rgb = RGBColor(236, 240, 241)
    
    # Title
    title_box2 = slide2.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame2 = title_box2.text_frame
    title_frame2.text = "¿Qué son los Modelos 3D?"
    title_para2 = title_frame2.paragraphs[0]
    title_para2.font.size = Pt(44)
    title_para2.font.bold = True
    title_para2.font.color.rgb = RGBColor(44, 62, 80)
    
    # Content
    content_box2 = slide2.shapes.add_textbox(Inches(1), Inches(1.8), Inches(8), Inches(4.5))
    content_frame2 = content_box2.text_frame
    content_frame2.word_wrap = True
    
    p1 = content_frame2.paragraphs[0]
    p1.text = "• Representaciones digitales de objetos tridimensionales"
    p1.font.size = Pt(24)
    p1.font.color.rgb = RGBColor(52, 73, 94)
    p1.space_after = Pt(20)
    
    p2 = content_frame2.add_paragraph()
    p2.text = "• Utilizados en diseño, arquitectura, videojuegos y animación"
    p2.font.size = Pt(24)
    p2.font.color.rgb = RGBColor(52, 73, 94)
    p2.space_after = Pt(20)
    
    p3 = content_frame2.add_paragraph()
    p3.text = "• Permiten visualización interactiva desde múltiples ángulos"
    p3.font.size = Pt(24)
    p3.font.color.rgb = RGBColor(52, 73, 94)
    p3.space_after = Pt(20)
    
    p4 = content_frame2.add_paragraph()
    p4.text = "• Formatos comunes: OBJ, FBX, GLB, GLTF, STL"
    p4.font.size = Pt(24)
    p4.font.color.rgb = RGBColor(52, 73, 94)
    
    # Slide 3: Cube Model
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])
    background3 = slide3.background
    fill3 = background3.fill
    fill3.solid()
    fill3.fore_color.rgb = RGBColor(236, 240, 241)
    
    # Title
    title_box3 = slide3.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame3 = title_box3.text_frame
    title_frame3.text = "Modelo 3D: Cubo"
    title_para3 = title_frame3.paragraphs[0]
    title_para3.font.size = Pt(44)
    title_para3.font.bold = True
    title_para3.font.color.rgb = RGBColor(44, 62, 80)
    
    # 3D Model placeholder
    model_box3 = slide3.shapes.add_shape(
        1,  # Rectangle
        Inches(2), Inches(2), Inches(6), Inches(4)
    )
    model_box3.fill.solid()
    model_box3.fill.fore_color.rgb = RGBColor(52, 152, 219)
    model_box3.line.color.rgb = RGBColor(41, 128, 185)
    model_box3.line.width = Pt(3)
    
    # Text in placeholder
    text_frame3 = model_box3.text_frame
    text_frame3.text = "🎲\n\nModelo 3D: CUBO\n\nGeometría básica con 6 caras\n8 vértices y 12 aristas"
    text_frame3.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame3.paragraphs[0].font.size = Pt(20)
    text_frame3.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    text_frame3.vertical_anchor = 1  # Middle
    
    # Slide 4: Sphere Model
    slide4 = prs.slides.add_slide(prs.slide_layouts[6])
    background4 = slide4.background
    fill4 = background4.fill
    fill4.solid()
    fill4.fore_color.rgb = RGBColor(236, 240, 241)
    
    # Title
    title_box4 = slide4.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame4 = title_box4.text_frame
    title_frame4.text = "Modelo 3D: Esfera"
    title_para4 = title_frame4.paragraphs[0]
    title_para4.font.size = Pt(44)
    title_para4.font.bold = True
    title_para4.font.color.rgb = RGBColor(44, 62, 80)
    
    # 3D Model placeholder
    model_box4 = slide4.shapes.add_shape(
        9,  # Oval
        Inches(2.5), Inches(2), Inches(5), Inches(4)
    )
    model_box4.fill.solid()
    model_box4.fill.fore_color.rgb = RGBColor(231, 76, 60)
    model_box4.line.color.rgb = RGBColor(192, 57, 43)
    model_box4.line.width = Pt(3)
    
    # Text in placeholder
    text_frame4 = model_box4.text_frame
    text_frame4.text = "⚽\n\nModelo 3D: ESFERA\n\nSuperficie curva perfecta\nTodos los puntos equidistantes del centro"
    text_frame4.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame4.paragraphs[0].font.size = Pt(20)
    text_frame4.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    text_frame4.vertical_anchor = 1
    
    # Slide 5: Pyramid Model
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])
    background5 = slide5.background
    fill5 = background5.fill
    fill5.solid()
    fill5.fore_color.rgb = RGBColor(236, 240, 241)
    
    # Title
    title_box5 = slide5.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame5 = title_box5.text_frame
    title_frame5.text = "Modelo 3D: Pirámide"
    title_para5 = title_frame5.paragraphs[0]
    title_para5.font.size = Pt(44)
    title_para5.font.bold = True
    title_para5.font.color.rgb = RGBColor(44, 62, 80)
    
    # 3D Model placeholder (triangle)
    model_box5 = slide5.shapes.add_shape(
        3,  # Triangle
        Inches(2.5), Inches(2), Inches(5), Inches(4)
    )
    model_box5.fill.solid()
    model_box5.fill.fore_color.rgb = RGBColor(241, 196, 15)
    model_box5.line.color.rgb = RGBColor(243, 156, 18)
    model_box5.line.width = Pt(3)
    
    # Text in placeholder
    text_frame5 = model_box5.text_frame
    text_frame5.text = "🔺\n\nPIRÁMIDE\n\nBase poligonal\nCara triangular"
    text_frame5.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame5.paragraphs[0].font.size = Pt(20)
    text_frame5.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    text_frame5.vertical_anchor = 1
    
    # Slide 6: Applications
    slide6 = prs.slides.add_slide(prs.slide_layouts[6])
    background6 = slide6.background
    fill6 = background6.fill
    fill6.solid()
    fill6.fore_color.rgb = RGBColor(236, 240, 241)
    
    # Title
    title_box6 = slide6.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
    title_frame6 = title_box6.text_frame
    title_frame6.text = "Aplicaciones de Modelos 3D"
    title_para6 = title_frame6.paragraphs[0]
    title_para6.font.size = Pt(44)
    title_para6.font.bold = True
    title_para6.font.color.rgb = RGBColor(44, 62, 80)
    
    # Content boxes
    # Box 1
    box1 = slide6.shapes.add_shape(1, Inches(0.8), Inches(2), Inches(4), Inches(2))
    box1.fill.solid()
    box1.fill.fore_color.rgb = RGBColor(46, 204, 113)
    box1.line.color.rgb = RGBColor(39, 174, 96)
    box1.line.width = Pt(2)
    text1 = box1.text_frame
    text1.text = "🎮 Videojuegos\n\nPersonajes, escenarios\ny objetos interactivos"
    text1.paragraphs[0].alignment = PP_ALIGN.CENTER
    text1.paragraphs[0].font.size = Pt(18)
    text1.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    text1.vertical_anchor = 1
    
    # Box 2
    box2 = slide6.shapes.add_shape(1, Inches(5.2), Inches(2), Inches(4), Inches(2))
    box2.fill.solid()
    box2.fill.fore_color.rgb = RGBColor(155, 89, 182)
    box2.line.color.rgb = RGBColor(142, 68, 173)
    box2.line.width = Pt(2)
    text2 = box2.text_frame
    text2.text = "🏗️ Arquitectura\n\nDiseño de edificios\ny visualización"
    text2.paragraphs[0].alignment = PP_ALIGN.CENTER
    text2.paragraphs[0].font.size = Pt(18)
    text2.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    text2.vertical_anchor = 1
    
    # Box 3
    box3 = slide6.shapes.add_shape(1, Inches(0.8), Inches(4.5), Inches(4), Inches(2))
    box3.fill.solid()
    box3.fill.fore_color.rgb = RGBColor(52, 152, 219)
    box3.line.color.rgb = RGBColor(41, 128, 185)
    box3.line.width = Pt(2)
    text3 = box3.text_frame
    text3.text = "🎬 Animación\n\nPelículas y efectos\nespeciales"
    text3.paragraphs[0].alignment = PP_ALIGN.CENTER
    text3.paragraphs[0].font.size = Pt(18)
    text3.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    text3.vertical_anchor = 1
    
    # Box 4
    box4 = slide6.shapes.add_shape(1, Inches(5.2), Inches(4.5), Inches(4), Inches(2))
    box4.fill.solid()
    box4.fill.fore_color.rgb = RGBColor(230, 126, 34)
    box4.line.color.rgb = RGBColor(211, 84, 0)
    box4.line.width = Pt(2)
    text4 = box4.text_frame
    text4.text = "🖨️ Impresión 3D\n\nPrototipos y\nfabricación"
    text4.paragraphs[0].alignment = PP_ALIGN.CENTER
    text4.paragraphs[0].font.size = Pt(18)
    text4.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    text4.vertical_anchor = 1
    
    # Slide 7: Conclusion
    slide7 = prs.slides.add_slide(prs.slide_layouts[6])
    background7 = slide7.background
    fill7 = background7.fill
    fill7.solid()
    fill7.fore_color.rgb = RGBColor(44, 62, 80)
    
    # Title
    title_box7 = slide7.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(1))
    title_frame7 = title_box7.text_frame
    title_frame7.text = "¡Gracias!"
    title_para7 = title_frame7.paragraphs[0]
    title_para7.font.size = Pt(60)
    title_para7.font.bold = True
    title_para7.font.color.rgb = RGBColor(255, 255, 255)
    title_para7.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    subtitle_box7 = slide7.shapes.add_textbox(Inches(1), Inches(4), Inches(8), Inches(1))
    subtitle_frame7 = subtitle_box7.text_frame
    subtitle_frame7.text = "El futuro es tridimensional"
    subtitle_para7 = subtitle_frame7.paragraphs[0]
    subtitle_para7.font.size = Pt(32)
    subtitle_para7.font.color.rgb = RGBColor(236, 240, 241)
    subtitle_para7.alignment = PP_ALIGN.CENTER
    
    # Save presentation
    output_file = '/vercel/sandbox/Presentacion_Modelos_3D.pptx'
    prs.save(output_file)
    print(f"✅ Presentación creada exitosamente: {output_file}")
    print(f"📊 Total de diapositivas: {len(prs.slides)}")
    
    return output_file

if __name__ == "__main__":
    create_3d_presentation()
